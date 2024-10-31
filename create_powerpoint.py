from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

title_slide_index = 0
image_slide_index = 8

def add_image_slide(presentation, title, image_path):
    slide_with_title_and_image_layout = presentation.slide_layouts[image_slide_index]
    slide = presentation.slides.add_slide(slide_with_title_and_image_layout)
    slide.shapes.title.text = title
    picture_placeholder = slide.placeholders[1]  # keyed by idx, not offset
    picture_placeholder.insert_picture(image_path)


def add_title_slide(presentation):
    layout = presentation.slide_layouts[title_slide_index]
    slide = presentation.slides.add_slide(layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]  # placeholder for subtitle
    title.text = "What if we could automate business validation tasks"  # title
    subtitle.text = "Imagine paging through a slide deck rather than clicking through all the sites"  # subtitle

def create_presentation_from_successful_test(pathPrefix):
    fileLocation = 'output/' + pathPrefix + "/"
    presentation = create_presentation()
    add_title_slide(presentation)
    add_image_slide(presentation, "b stack demo", fileLocation + "username.png")
    add_image_slide(presentation, "amazon", fileLocation + "amazon.png")

    presentation.save('output/' + pathPrefix + ".pptx") 

def create_presentation():
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)
    return presentation # saving file


if __name__ == "__main__":
    savedFileDirectory = sys.argv[1]
    create_presentation_from_successful_test(savedFileDirectory)
